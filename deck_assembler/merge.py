"""Mixed-deck merge: image-PPTX + chart-PPTX → single editable PPTX, slides in NN order.

Strategy:
 1. Group every artifact in the deck dir into (NN, kind, path), where kind is
    'image' for *.png and 'chart' for *.svg slide files.
 2. Run editable_pptx over the image subset → tmp/image-deck.pptx.
 3. Run svg_to_pptx over the chart subset → tmp/chart-deck.pptx.
 4. Open both source decks. Build the destination deck by appending slides in
    NN order, copying each <p:sld> XML and remapping its image relationships.
"""

from __future__ import annotations

import logging
import re
import shutil
import tempfile
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

logger = logging.getLogger(__name__)

SLIDE_PATTERN = re.compile(r"^(\d+)-slide-.*\.(png|jpg|jpeg|svg)$", re.IGNORECASE)


def _list_slide_artifacts(deck_dir: Path) -> list[tuple[int, str, Path]]:
    """Return [(NN, kind, path)] sorted by NN. kind ∈ {'image','chart'}."""
    out: list[tuple[int, str, Path]] = []
    for f in deck_dir.iterdir():
        if not f.is_file():
            continue
        m = SLIDE_PATTERN.match(f.name)
        if not m:
            continue
        nn = int(m.group(1))
        ext = m.group(2).lower()
        kind = "chart" if ext == "svg" else "image"
        out.append((nn, kind, f))
    out.sort(key=lambda t: t[0])
    return out


def _build_image_deck(slides: list[Path], out_pptx: Path, deck_dir: Path) -> None:
    """Run the existing editable_pptx pipeline on the supplied PNG slides."""
    from editable_pptx.assemble import export_editable_deck
    from editable_pptx.canvas import materialize_normalized_slides, resolve_target_canvas_wh
    from editable_pptx.env import (
        background_mode,
        load_skilldeck_env,
        mineru_config,
        mineru_poll_timeout,
    )
    from editable_pptx.mineru import MinerUError, parse_slide_image

    repo_root = Path(__file__).resolve().parent.parent
    load_skilldeck_env(repo_root)
    cfg = mineru_config()
    if not cfg["token"]:
        raise RuntimeError(
            "MINERU_TOKEN missing — required for image-slide editable export. "
            "Set it in `.env` or use a chart-only deck."
        )

    bg_mode = background_mode()
    timeout = mineru_poll_timeout()

    with tempfile.TemporaryDirectory(prefix="deck_assembler_img_") as work:
        work_root = Path(work)
        tw, th = resolve_target_canvas_wh(slides)
        normalized_dir = work_root / "normalized_slides"
        slides_for_pipeline = materialize_normalized_slides(slides, normalized_dir, tw, th)

        mineru_dirs: list[Path] = []
        for i, slide_path in enumerate(slides_for_pipeline):
            wd = work_root / f"slide_{i:03d}"
            logger.info("MinerU parse %s/%s: %s", i + 1, len(slides), slide_path.name)
            try:
                mdir = parse_slide_image(
                    str(slide_path),
                    token=cfg["token"],
                    api_base=cfg["api_base"],
                    model_version=cfg["model_version"],
                    work_dir=wd,
                    poll_timeout=timeout,
                )
                mineru_dirs.append(mdir)
            except MinerUError as e:
                raise RuntimeError(f"MinerU failed for {slide_path.name}: {e}") from e

        export_editable_deck(
            [str(s) for s in slides_for_pipeline],
            mineru_dirs,
            out_pptx,
            bg_mode=bg_mode,
            deck_dir=deck_dir,
        )


def _build_chart_deck(svg_paths: list[Path], out_pptx: Path) -> None:
    """Run svg_to_pptx on the chart SVGs.

    use_native_shapes=True turns each SVG primitive into a real DrawingML shape,
    which is the whole point of the ppt-master export pipeline. Animation and
    notes are disabled here because the hybrid deck mixes content kinds and
    those features assume a uniform slide structure.
    """
    from svg_to_pptx.pptx_builder import create_pptx_with_native_svg

    create_pptx_with_native_svg(
        svg_files=svg_paths,
        output_path=out_pptx,
        canvas_format="ppt169",
        verbose=False,
        animation=None,
        animation_config=None,
        enable_notes=False,
        use_native_shapes=True,
    )


_SLIDE_PART_RELS_TO_SKIP = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
)

_OFF_TAG = qn("a:off")
_EXT_TAG = qn("a:ext")
_CHEXT_TAG = qn("a:chExt")
_CHOFF_TAG = qn("a:chOff")


def _scale_spTree(spTree, scale_x: float, scale_y: float) -> None:
    """Scale every absolute coordinate inside spTree by (scale_x, scale_y).

    Walks every <a:off> / <a:ext> (and the <a:chOff>/<a:chExt> inside group
    transforms) and rewrites its x/y or cx/cy. This is a flat coordinate scale
    — adequate when the source and destination both use the same orientation
    and only differ in canvas size, which is the case when a chart deck (default
    ppt169) is appended onto a wider image deck.
    """
    if abs(scale_x - 1.0) < 1e-6 and abs(scale_y - 1.0) < 1e-6:
        return
    for elem in spTree.iter():
        tag = elem.tag
        if tag in (_OFF_TAG, _CHOFF_TAG):
            x = elem.get("x")
            y = elem.get("y")
            if x is not None:
                elem.set("x", str(int(round(int(x) * scale_x))))
            if y is not None:
                elem.set("y", str(int(round(int(y) * scale_y))))
        elif tag in (_EXT_TAG, _CHEXT_TAG):
            cx = elem.get("cx")
            cy = elem.get("cy")
            if cx is not None:
                elem.set("cx", str(int(round(int(cx) * scale_x))))
            if cy is not None:
                elem.set("cy", str(int(round(int(cy) * scale_y))))


def _copy_slide(
    dest_prs: Presentation,
    source_slide,
    *,
    source_w: int | None = None,
    source_h: int | None = None,
) -> None:
    """Append a deep copy of source_slide into dest_prs.

    If the source presentation's slide dimensions differ from dest_prs, every
    coordinate inside the copied spTree is rescaled so the slide fills the dest
    canvas instead of clustering in the upper-left corner.
    """
    blank = _pick_blank_layout(dest_prs)
    new_slide = dest_prs.slides.add_slide(blank)

    new_spTree = deepcopy(source_slide.shapes._spTree)
    parent = new_slide.shapes._spTree.getparent()
    parent.replace(new_slide.shapes._spTree, new_spTree)

    if source_w and source_h:
        sx = dest_prs.slide_width / source_w
        sy = dest_prs.slide_height / source_h
        _scale_spTree(new_spTree, sx, sy)

    for rel in source_slide.part.rels.values():
        if rel.is_external or rel.reltype in _SLIDE_PART_RELS_TO_SKIP:
            continue
        new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        _rewrite_relationship_ids(new_spTree, rel.rId, new_rid)


def _pick_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in (layout.name or "").lower():
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


_REL_ATTRS = (
    qn("r:embed"),
    qn("r:link"),
    qn("r:id"),
)


def _rewrite_relationship_ids(spTree, old_rid: str, new_rid: str) -> None:
    """Walk spTree, replacing every relationship-reference attribute that holds old_rid."""
    if old_rid == new_rid:
        return
    for elem in spTree.iter():
        for attr in _REL_ATTRS:
            if elem.get(attr) == old_rid:
                elem.set(attr, new_rid)


def _set_dimensions_from(prs: Presentation, source_prs: Presentation) -> None:
    prs.slide_width = source_prs.slide_width
    prs.slide_height = source_prs.slide_height


def assemble_mixed_deck(deck_dir: Path, output_pptx: Path) -> None:
    artifacts = _list_slide_artifacts(deck_dir)
    if not artifacts:
        raise SystemExit(f"No slide artifacts (NN-slide-*.png|svg) in {deck_dir}")

    image_slides = [p for nn, k, p in artifacts if k == "image"]
    chart_slides = [p for nn, k, p in artifacts if k == "chart"]

    logger.info("Found %d image slide(s), %d chart slide(s)", len(image_slides), len(chart_slides))

    with tempfile.TemporaryDirectory(prefix="deck_assembler_") as work:
        work_root = Path(work)
        image_pptx = work_root / "image-deck.pptx"
        chart_pptx = work_root / "chart-deck.pptx"

        if image_slides:
            _build_image_deck(image_slides, image_pptx, deck_dir)
        if chart_slides:
            _build_chart_deck(chart_slides, chart_pptx)

        # Open source presentations.
        img_prs = Presentation(str(image_pptx)) if image_slides else None
        cht_prs = Presentation(str(chart_pptx)) if chart_slides else None

        # Build destination — clone the dimensions from the first available source so
        # both content kinds get the same canvas. Image side wins when both exist
        # (chart slides scale to fill via their viewBox).
        dest = Presentation()
        if img_prs is not None:
            _set_dimensions_from(dest, img_prs)
        elif cht_prs is not None:
            _set_dimensions_from(dest, cht_prs)

        # Map filename → (source_prs, source_slide)
        per_slide: dict[Path, tuple[Presentation, object]] = {}
        if img_prs is not None:
            for path, slide in zip(image_slides, list(img_prs.slides)):
                per_slide[path] = (img_prs, slide)
        if cht_prs is not None:
            for path, slide in zip(chart_slides, list(cht_prs.slides)):
                per_slide[path] = (cht_prs, slide)

        for nn, kind, path in artifacts:
            mapping = per_slide.get(path)
            if mapping is None:
                logger.warning("No source slide for %s — skipped", path.name)
                continue
            src_prs, src_slide = mapping
            logger.info("Append slide %02d (%s): %s", nn, kind, path.name)
            _copy_slide(
                dest,
                src_slide,
                source_w=src_prs.slide_width,
                source_h=src_prs.slide_height,
            )

        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        dest.core_properties.title = output_pptx.stem
        dest.save(str(output_pptx))
