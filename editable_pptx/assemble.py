"""Build editable PPTX with python-pptx (background + native shapes + pictures + text)."""

from __future__ import annotations

import logging
import tempfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from editable_pptx.background import build_background
from editable_pptx.fonts import (
    WEIGHT_TO_BOLD,
    fit_font_size_pt,
    resolve_family,
    text_has_cjk,
)
from editable_pptx.layout import RASTER_TYPES

logger = logging.getLogger(__name__)

DEFAULT_DPI = 96


def _pixels_to_inches(px: float) -> float:
    return px / DEFAULT_DPI


def _px_to_emu(px: float) -> int:
    return int(round(px / DEFAULT_DPI * 914400))


def _pp_align(name: str | None) -> int:
    n = (name or "left").lower()
    if n == "center":
        return PP_ALIGN.CENTER
    if n == "right":
        return PP_ALIGN.RIGHT
    if n == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


_KIND_TO_MSO: dict[str, int] = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
    "line": MSO_SHAPE.RECTANGLE,  # rendered as a thin filled rect; LINE shape needs end-points
    "diamond": MSO_SHAPE.DIAMOND,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
}


def _add_native_shape(slide, shape: dict[str, Any]):
    bb = shape["bbox"]
    x0, y0, x1, y1 = bb
    w = max(1.0, x1 - x0)
    h = max(1.0, y1 - y0)
    kind = shape["kind"]

    # `line` becomes a 1-2px filled rectangle so we keep one consistent path.
    if kind == "line":
        thickness = max(1.0, shape.get("stroke_width_px") or 2.0)
        if w >= h:
            h = thickness
        else:
            w = thickness

    mso = _KIND_TO_MSO.get(kind, MSO_SHAPE.RECTANGLE)
    sp = slide.shapes.add_shape(
        mso,
        Emu(_px_to_emu(x0)),
        Emu(_px_to_emu(y0)),
        Emu(_px_to_emu(w)),
        Emu(_px_to_emu(h)),
    )

    # Corner radius for round-rects and pills via the first adjustment handle.
    if kind in ("roundRect", "pill"):
        try:
            short_side = min(w, h)
            if kind == "pill":
                adj = 0.5
            else:
                cr = shape.get("corner_radius_px") or 0
                adj = max(0.0, min(0.5, (cr / short_side) if short_side else 0))
                if adj < 0.05 and (cr or 0) > 0:
                    adj = 0.1
            if sp.adjustments:
                sp.adjustments[0] = adj
        except Exception:
            pass

    fill_rgb = shape.get("fill_rgb")
    stroke_rgb = shape.get("stroke_rgb")
    stroke_w_px = shape.get("stroke_width_px") or 0
    try:
        if fill_rgb:
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(*fill_rgb)
        else:
            sp.fill.background()
    except Exception:
        pass
    try:
        if stroke_rgb:
            sp.line.color.rgb = RGBColor(*stroke_rgb)
            if stroke_w_px:
                sp.line.width = Emu(_px_to_emu(stroke_w_px))
        elif stroke_w_px == 0 and not fill_rgb:
            # neither fill nor stroke -> at least keep a faint outline so user can find it
            pass
        else:
            sp.line.fill.background()
    except Exception:
        pass

    try:
        if sp.has_text_frame:
            sp.text_frame.text = ""
    except Exception:
        pass
    return sp


def _add_flat_background(slide, slide_w_px: int, slide_h_px: int, page_bg: tuple[int, int, int] | None) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Emu(_px_to_emu(slide_w_px)),
        Emu(_px_to_emu(slide_h_px)),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*(page_bg or (255, 255, 255)))
    bg.line.fill.background()


def add_slide_from_image(
    prs: Presentation,
    slide_image_path: str,
    elements: list[dict[str, Any]],
    *,
    bg_mode: str,
    set_presentation_dimensions: bool,
    text_pad: float = 1.005,
    font_body: str | None = None,
    font_title: str | None = None,
    shapes: list[dict[str, Any]] | None = None,
    page_bg: tuple[int, int, int] | None = None,
    flat_background: bool = False,
) -> None:
    """
    Add one slide. python-pptx uses one slide size for the whole deck — set
    `set_presentation_dimensions=True` only for the first slide (all PNGs should match).
    """
    im = Image.open(slide_image_path)
    slide_w_px, slide_h_px = im.size
    if set_presentation_dimensions:
        prs.slide_width = Inches(_pixels_to_inches(slide_w_px))
        prs.slide_height = Inches(_pixels_to_inches(slide_h_px))

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    shapes = shapes or []
    grouped_refs: list[tuple[str | None, str | None, Any]] = []

    # --- background layer ---
    if flat_background:
        _add_flat_background(slide, slide_w_px, slide_h_px, page_bg)
    else:
        bg_img = build_background(slide_image_path, elements, mode=bg_mode, shapes=shapes)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = Path(tmp.name)
            bg_img.save(tmp_path, "PNG")
        try:
            slide.shapes.add_picture(
                str(tmp_path),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        finally:
            tmp_path.unlink(missing_ok=True)

    # --- native shapes that go UNDER text/pictures ---
    for sh in shapes:
        if sh.get("z", "under") == "under":
            try:
                sp = _add_native_shape(slide, sh)
                grouped_refs.append((sh.get("parent_id"), sh.get("candidate_id"), sp))
            except Exception as e:
                logger.warning("under shape %s failed: %s", sh.get("kind"), e)

    # --- pictures (images / tables / charts from MinerU) ---
    pictures: list[dict[str, Any]] = []
    texts: list[dict[str, Any]] = []
    for el in elements:
        t = el.get("type", "")
        if t in RASTER_TYPES and el.get("image_path"):
            pictures.append(el)
        elif el.get("content"):
            texts.append(el)

    for el in pictures:
        path = el["image_path"]
        bb = el["bbox"]
        if not path or not Path(path).is_file():
            continue
        left = Inches(_pixels_to_inches(bb[0]))
        top = Inches(_pixels_to_inches(bb[1]))
        w = Inches(_pixels_to_inches(bb[2] - bb[0]))
        h = Inches(_pixels_to_inches(bb[3] - bb[1]))
        try:
            slide.shapes.add_picture(path, left=left, top=top, width=w, height=h)
        except Exception as e:
            logger.warning("add_picture failed: %s", e)

    # --- text boxes ---
    for el in texts:
        text = el.get("content") or ""
        if not text.strip():
            continue
        bb = el["bbox"]
        pad = text_pad
        cx = (bb[0] + bb[2]) / 2
        cy = (bb[1] + bb[3]) / 2
        bw = (bb[2] - bb[0]) * pad
        bh = (bb[3] - bb[1]) * pad
        x0 = cx - bw / 2
        y0 = cy - bh / 2
        left = Inches(_pixels_to_inches(x0))
        top = Inches(_pixels_to_inches(y0))
        width = Inches(_pixels_to_inches(bw))
        height = Inches(_pixels_to_inches(bh))
        box = slide.shapes.add_textbox(left, top, width, height)
        grouped_refs.append((el.get("parent_id"), None, box))
        try:
            box.fill.background()
        except Exception:
            pass
        try:
            box.line.fill.background()
        except Exception:
            pass
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        try:
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.text = text
        sty = el.get("style") or {}
        p.alignment = _pp_align(sty.get("align"))
        try:
            p.line_spacing = 1.08
        except Exception:
            pass

        is_title = el.get("type") == "title"
        # Pick font family.
        family_override = font_title if is_title and font_title else font_body
        family = resolve_family(
            sty.get("font_family_hint"),
            text_has_cjk(text),
            override=family_override,
        )
        # Shrink-to-fit using the chosen family.
        fs = fit_font_size_pt(text, bw, bh, family)
        # Bold from weight if VLM gave one, else from `bold` flag, else title heuristic.
        weight = (sty.get("weight") or "regular").lower()
        bold_from_weight = WEIGHT_TO_BOLD.get(weight)
        bold = sty.get("bold")
        if bold is None:
            bold = bold_from_weight if bold_from_weight is not None else is_title
        italic = bool(sty.get("italic", False))
        underline = bool(sty.get("underline", False))
        rgb = sty.get("color_rgb")

        for run in p.runs:
            run.font.size = Pt(max(6, int(round(fs))))
            run.font.bold = bool(bold)
            run.font.italic = italic
            if underline:
                run.font.underline = True
            if rgb and isinstance(rgb, (list, tuple)) and len(rgb) >= 3:
                try:
                    run.font.color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))
                except Exception:
                    pass
            try:
                run.font.name = family
            except Exception:
                pass

    # --- native shapes that go OVER text ---
    for sh in shapes:
        if sh.get("z", "under") == "over":
            try:
                sp = _add_native_shape(slide, sh)
                grouped_refs.append((sh.get("parent_id"), sh.get("candidate_id"), sp))
            except Exception as e:
                logger.warning("over shape %s failed: %s", sh.get("kind"), e)

    _group_related_shapes(slide, grouped_refs)


def _group_related_shapes(slide, refs: list[tuple[str | None, str | None, Any]]) -> None:
    """Group container shapes with direct child shapes/text boxes when ids are available."""
    by_candidate: dict[str, Any] = {cid: shape for _pid, cid, shape in refs if cid}
    by_parent: dict[str, list[Any]] = {}
    for pid, cid, shape in refs:
        if not pid or pid == cid:
            continue
        by_parent.setdefault(pid, []).append(shape)
    # Inner groups first, then outer groups. This is intentionally conservative:
    # only direct parent relationships from the CV hierarchy are grouped.
    for parent_id, children in sorted(by_parent.items(), key=lambda item: len(item[1])):
        parent_shape = by_candidate.get(parent_id)
        members = ([parent_shape] if parent_shape is not None else []) + children
        if len(members) < 2:
            continue
        try:
            slide.shapes.add_group_shape(members)
        except Exception as e:
            logger.debug("group shape failed for %s: %s", parent_id, e)


def export_editable_deck(
    ordered_image_paths: list[str],
    mineru_dirs: list[Path | None],
    output_pptx: Path,
    *,
    bg_mode: str,
    deck_dir: Path | None = None,
) -> None:
    if len(ordered_image_paths) != len(mineru_dirs):
        raise ValueError("image paths and mineru_dirs length mismatch")
    if not ordered_image_paths:
        raise ValueError("no slides")

    w0, h0 = Image.open(ordered_image_paths[0]).size
    for p in ordered_image_paths[1:]:
        w, h = Image.open(p).size
        if (w, h) != (w0, h0):
            logger.warning(
                "Slide size mismatch: first is %sx%s but %s is %sx%s — PPT uses first size; positions may drift.",
                w0,
                h0,
                p,
                w,
                h,
            )

    from editable_pptx.env import (
        analysis_cache_enabled,
        analysis_workers,
        bg_flatten_enabled,
        crosscheck_enabled,
        diagram_decompose_enabled,
        diagram_decompose_max_items,
        diagram_decompose_min_area_fraction,
        font_config,
        hybrid_cv_enabled,
        hybrid_min_area_fraction,
        hybrid_recursion_depth,
        layout_engine,
        layout_snap_enabled,
        shape_detect_enabled,
        snap_cluster_tol_px,
        snap_grid_px,
        text_pad_ratio,
        vlm_enabled,
        vlm_style_model,
    )
    from editable_pptx import analysis_cache
    from editable_pptx.decompose import decompose_image_regions
    from editable_pptx.layout import elements_from_mineru_dir
    from editable_pptx.openai_style import (
        apply_styles_to_elements,
        compute_openai_element_styles,
    )
    from editable_pptx.shapes import detect_shapes
    from editable_pptx.snap import edge_snap_bboxes, snap_bboxes

    fc = font_config(deck_dir)
    pad = text_pad_ratio()
    do_snap = layout_snap_enabled()
    do_shapes = shape_detect_enabled()
    do_flat_bg = bg_flatten_enabled()
    do_hybrid = hybrid_cv_enabled()
    do_decompose = diagram_decompose_enabled()
    do_cache = analysis_cache_enabled()
    grid = snap_grid_px()
    cluster_tol = snap_cluster_tol_px()
    decompose_min_area = diagram_decompose_min_area_fraction()
    decompose_max_items = diagram_decompose_max_items()
    vlm_model_name = vlm_style_model()
    layout_engine_name = layout_engine()

    def analyze_one(idx: int, img: str, mdir: Path | None) -> dict[str, Any]:
        """Run layout + VLM + decompose for one slide.

        Returns a dict with the per-slide state needed by the serial assembly
        phase: `idx`, `img`, `slide_size`, `els`, `shapes_list`, `page_bg`.
        """
        im = Image.open(img)
        slide_size = im.size
        page_bg: tuple[int, int, int] | None = None
        shapes_list: list[dict[str, Any]] = []
        if do_hybrid:
            from editable_pptx.hybrid import analyze_slide_hybrid

            els, shapes_list, page_bg, debug = analyze_slide_hybrid(
                img,
                mineru_dir=mdir,
                min_area_fraction=hybrid_min_area_fraction(),
                recursion_depth=hybrid_recursion_depth(),
            )
            logger.info("Hybrid layout analysis for %s: %s", img, debug)
        else:
            els = elements_from_mineru_dir(mdir, slide_size)

        slide_stem = Path(img).stem
        cache_file = analysis_cache.cache_path(deck_dir, slide_stem) if do_cache else None
        cache_hit: dict[str, Any] | None = None
        key: str | None = None
        if do_cache and vlm_enabled():
            key = analysis_cache.cache_key(
                image_path=img,
                vlm_model=vlm_model_name,
                decompose_enabled=do_decompose,
                decompose_min_area_fraction=decompose_min_area,
                layout_engine=layout_engine_name,
            )
            cache_hit = analysis_cache.load(cache_file, key)

        if cache_hit is not None:
            logger.info("Analysis cache HIT for %s", img)
            cached_styles = cache_hit.get("styles") or []
            apply_styles_to_elements(els, cached_styles)
            cached_bg = cache_hit.get("page_bg")
            if cached_bg:
                page_bg = page_bg or tuple(cached_bg)
            shapes_list.extend(cache_hit.get("shapes") or [])
            decompose_block = cache_hit.get("decompose") or {}
            extra_shapes = decompose_block.get("extra_shapes") or []
            extra_texts = decompose_block.get("extra_texts") or []
            removed_idx = decompose_block.get("removed_indices") or []
            if removed_idx:
                drop = set(removed_idx)
                els = [e for i, e in enumerate(els) if i not in drop]
            if extra_texts:
                els.extend(extra_texts)
            if extra_shapes:
                shapes_list.extend(extra_shapes)
        else:
            captured_styles: list[dict[str, Any]] = []
            captured_bg: tuple[int, int, int] | None = None
            if vlm_enabled():
                logger.info("VLM style extraction for %s", img)
                captured_styles, captured_bg = compute_openai_element_styles(img, els)
                if captured_styles:
                    apply_styles_to_elements(els, captured_styles)
                page_bg = page_bg or captured_bg

            layout_hints: list[str] = []
            spec = _read_spec_for(img, deck_dir)
            if spec:
                arch = spec.get("layout")
                if arch:
                    layout_hints.append(f"layout={arch}")
                for sh in spec.get("shape_hints") or []:
                    k = sh.get("kind")
                    if k:
                        layout_hints.append(f"hint:{k}")

            captured_shapes: list[dict[str, Any]] = []
            if do_shapes and not do_hybrid:
                captured_shapes = detect_shapes(img, els, slide_size, layout_hints=layout_hints)
                shapes_list.extend(captured_shapes)
            if do_hybrid:
                edge_snap_bboxes(img, els, only_sources={"vlm_missing"})
                edge_snap_bboxes(img, shapes_list, only_sources={"vlm_missing"})

            extra_shapes_d: list[dict[str, Any]] = []
            extra_texts_d: list[dict[str, Any]] = []
            removed_idx_d: list[int] = []
            if do_decompose:
                extra_shapes_d, extra_texts_d, removed_idx_d = decompose_image_regions(
                    img,
                    els,
                    slide_size,
                    min_area_fraction=decompose_min_area,
                    max_items_per_region=decompose_max_items,
                )
                if removed_idx_d:
                    drop = set(removed_idx_d)
                    els = [e for i, e in enumerate(els) if i not in drop]
                if extra_texts_d:
                    els.extend(extra_texts_d)
                if extra_shapes_d:
                    shapes_list.extend(extra_shapes_d)

            if do_cache and vlm_enabled() and cache_file is not None and key is not None:
                analysis_cache.save(
                    cache_file,
                    key,
                    page_bg=captured_bg,
                    styles=captured_styles,
                    shapes=captured_shapes,
                    decompose_extra_shapes=extra_shapes_d,
                    decompose_extra_texts=extra_texts_d,
                    decompose_removed_indices=removed_idx_d,
                )

        if do_snap:
            snap_bboxes(
                els, grid_px=grid, cluster_tol_px=cluster_tol,
                slide_w_px=slide_size[0], slide_h_px=slide_size[1],
            )
            snap_bboxes(
                shapes_list, grid_px=grid, cluster_tol_px=cluster_tol,
                slide_w_px=slide_size[0], slide_h_px=slide_size[1],
            )

        return {
            "idx": idx,
            "img": img,
            "slide_size": slide_size,
            "els": els,
            "shapes_list": shapes_list,
            "page_bg": page_bg,
        }

    # ---- Phase 1: parallel analysis. Workers default to 4; cap at slide count.
    n_slides = len(ordered_image_paths)
    workers = min(analysis_workers(), n_slides)
    analysis_results: list[dict[str, Any]] = [None] * n_slides  # type: ignore[list-item]
    if workers > 1 and n_slides > 1:
        with ThreadPoolExecutor(max_workers=workers, thread_name_prefix="slide-analysis") as ex:
            futures = [
                ex.submit(analyze_one, i, img, mdir)
                for i, (img, mdir) in enumerate(zip(ordered_image_paths, mineru_dirs))
            ]
            for fut in futures:
                result = fut.result()
                analysis_results[result["idx"]] = result
    else:
        for i, (img, mdir) in enumerate(zip(ordered_image_paths, mineru_dirs)):
            analysis_results[i] = analyze_one(i, img, mdir)

    # ---- Phase 2: serial python-pptx assembly.
    prs = Presentation()
    first = True
    for r in analysis_results:
        add_slide_from_image(
            prs,
            r["img"],
            r["els"],
            bg_mode=bg_mode,
            set_presentation_dimensions=first,
            text_pad=pad,
            font_body=fc["body"],
            font_title=fc["title"],
            shapes=r["shapes_list"],
            page_bg=r["page_bg"],
            flat_background=do_flat_bg,
        )
        first = False

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.core_properties.title = output_pptx.stem
    prs.save(str(output_pptx))

    if crosscheck_enabled():
        try:
            from editable_pptx.crosscheck import crosscheck_deck

            logger.info("Crosscheck rendering via LibreOffice...")
            crosscheck_deck(output_pptx, [Path(p) for p in ordered_image_paths])
        except Exception as e:
            logger.warning("Crosscheck skipped: %s", e)


def _read_spec_for(image_path: str, deck_dir: Path | None = None) -> dict[str, Any] | None:
    """Look for `<deck_dir>/prompts/NN-slide-*.spec.json` (deck_dir defaults to image parent)."""
    import json
    import re as _re

    p = Path(image_path)
    m = _re.match(r"^(\d+)-slide-(.+)\.(png|jpg|jpeg)$", p.name, _re.IGNORECASE)
    if not m:
        return None
    base = deck_dir if deck_dir is not None else p.parent
    prompts_dir = base / "prompts"
    if not prompts_dir.is_dir():
        return None
    stem = f"{m.group(1)}-slide-{m.group(2)}"
    spec = prompts_dir / f"{stem}.spec.json"
    if not spec.is_file():
        return None
    try:
        return json.loads(spec.read_text(encoding="utf-8"))
    except (OSError, ValueError):
        return None
